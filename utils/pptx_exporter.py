from pathlib import Path
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from utils.text_cleaner import ensure_text, clean_ai_text, clean_heading, is_noise_line

DARK = RGBColor(7, 17, 31)
NAVY = RGBColor(15, 23, 42)
GREEN = RGBColor(0, 166, 118)
MINT = RGBColor(16, 185, 129)
BLUE = RGBColor(56, 189, 248)
WHITE = RGBColor(248, 250, 252)
MUTED = RGBColor(100, 116, 139)
LIGHT = RGBColor(241, 245, 249)
LINE = RGBColor(226, 232, 240)
TEXT = RGBColor(15, 23, 42)

def _clean_text(text: str) -> str:
    return clean_ai_text(text)

def _is_note(line: str) -> bool:
    s = _clean_text(line).lower()
    return s.startswith("catatan narasi") or s.startswith("narasi") or s.startswith("catatan presenter")

def _strip_note(line: str) -> str:
    s = _clean_text(line)
    return re.sub(r"^(Catatan narasi|Narasi|Catatan presenter)\s*:?\s*", "", s, flags=re.I).strip()

def _strip_slide_prefix(line: str) -> str:
    s = _clean_text(line)
    s = re.sub(r"^Slide\s+\d+\s*[:\-–]\s*", "", s, flags=re.I).strip()
    return s

def _strip_bullet(line: str) -> str:
    s = _clean_text(line)
    s = re.sub(r"^[-•]\s*", "", s).strip()
    s = re.sub(r"^\d+[\.\)]\s*", "", s).strip()
    return s

def parse_slides(outline_text):
    raw = clean_ai_text(outline_text)

    # Normalize heading like "### Slide 2 – Title" to "Slide 2: Title"
    raw = re.sub(r"(?im)^\s*#{1,6}\s*", "", raw)
    raw = re.sub(r"(?im)^(Slide\s+\d+)\s*[–-]\s*", r"\1: ", raw)

    pattern = re.compile(r"(?ims)^\s*Slide\s+(\d+)\s*:\s*(.*?)(?=^\s*Slide\s+\d+\s*:|\Z)")
    matches = list(pattern.finditer(raw))
    slides = []

    if matches:
        for m in matches:
            block = m.group(0).strip()
            lines = [clean_ai_text(l) for l in block.splitlines() if not is_noise_line(l)]
            if not lines:
                continue

            title = _strip_slide_prefix(lines[0])
            if _is_note(title):
                title = f"Slide {m.group(1)}"

            bullets = []
            notes = ""

            for line in lines[1:]:
                if is_noise_line(line):
                    continue
                if _is_note(line):
                    notes = _strip_note(line)
                    continue

                item = _strip_bullet(line)
                # Avoid adding repeated slide title or note content as bullet
                if not item or _is_note(item):
                    continue
                if re.match(r"^Slide\s+\d+", item, re.I):
                    item = _strip_slide_prefix(item)
                if item:
                    bullets.append(item)

            slides.append({"title": title, "bullets": bullets[:4], "notes": notes})

    if not slides:
        # fallback: build slides from sensible chunks, skipping notes as titles
        lines = [clean_ai_text(l) for l in raw.splitlines() if not is_noise_line(l)]
        lines = [l for l in lines if not _is_note(l)]
        if not lines:
            lines = ["Presentasi Layanan Zimeira Tech", "Data belum tersedia."]
        chunks, current = [], []
        for line in lines:
            if len(current) >= 5:
                chunks.append(current)
                current = []
            current.append(line)
        if current:
            chunks.append(current)
        for i, chunk in enumerate(chunks[:10], start=1):
            title = _strip_slide_prefix(chunk[0]) if chunk else f"Slide {i}"
            bullets = [_strip_bullet(x) for x in chunk[1:5] if _strip_bullet(x)]
            slides.append({"title": title, "bullets": bullets[:4], "notes": ""})

    # Remove accidental note-only slides and ensure bullets
    clean_slides = []
    for idx, s in enumerate(slides, start=1):
        title = _clean_text(s.get("title", f"Slide {idx}"))
        if _is_note(title):
            continue
        bullets = [_clean_text(b) for b in s.get("bullets", []) if _clean_text(b) and not _is_note(b)]
        if not bullets:
            bullets = ["Poin utama disesuaikan berdasarkan data calon pelanggan."]
        clean_slides.append({"title": title, "bullets": bullets[:4], "notes": _clean_text(s.get("notes", ""))})

    return clean_slides[:12]

def set_bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_top_brand(slide, dark=False):
    color = WHITE if dark else TEXT
    mark = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.38), Inches(0.72), Inches(0.42))
    mark.fill.solid()
    mark.fill.fore_color.rgb = GREEN
    mark.line.fill.background()
    tf = mark.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "ZT"
    r.font.bold = True
    r.font.size = Pt(12)
    r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    box = slide.shapes.add_textbox(Inches(1.42), Inches(0.40), Inches(4), Inches(0.32))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Zimeira Tech"
    r.font.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = color

def add_footer(slide, index=None, dark=False):
    color = RGBColor(203, 213, 225) if dark else MUTED
    left = slide.shapes.add_textbox(Inches(0.65), Inches(7.03), Inches(4.8), Inches(0.25))
    p = left.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Komputer • Desain • Website • Digital Service"
    r.font.size = Pt(9)
    r.font.color.rgb = color
    if index is not None:
        right = slide.shapes.add_textbox(Inches(11.8), Inches(7.03), Inches(0.8), Inches(0.25))
        p = right.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"{index:02d}"
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = color
        p.alignment = PP_ALIGN.RIGHT

def add_title_slide(prs, project_title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)

    a = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.9), Inches(-0.8), Inches(3.8), Inches(3.8))
    a.fill.solid(); a.fill.fore_color.rgb = GREEN; a.fill.transparency = 25; a.line.fill.background()
    b = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.9), Inches(4.2), Inches(2.4), Inches(2.4))
    b.fill.solid(); b.fill.fore_color.rgb = BLUE; b.fill.transparency = 45; b.line.fill.background()

    add_top_brand(slide, dark=True)

    kicker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(1.35), Inches(3.5), Inches(0.42))
    kicker.fill.solid(); kicker.fill.fore_color.rgb = RGBColor(10, 80, 67); kicker.line.color.rgb = RGBColor(22, 163, 127)
    tf = kicker.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "PRESENTASI PENAWARAN LAYANAN"
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = RGBColor(209, 250, 229)
    p.alignment = PP_ALIGN.CENTER

    title = slide.shapes.add_textbox(Inches(0.68), Inches(2.05), Inches(8.8), Inches(1.75))
    tf = title.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = _clean_text(project_title)
    r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.72), Inches(4.12), Inches(7.8), Inches(0.9))
    tf = sub.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "Disusun berdasarkan data calon pelanggan dan katalog layanan Zimeira Tech."
    r.font.size = Pt(17); r.font.color.rgb = RGBColor(203, 213, 225)

    add_footer(slide, dark=True)

def add_content_slide(prs, index, title_text, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT)
    add_top_brand(slide)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = GREEN; accent.line.fill.background()

    num = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.04), Inches(0.74), Inches(0.45))
    num.fill.solid(); num.fill.fore_color.rgb = RGBColor(220, 252, 231); num.line.fill.background()
    tf = num.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = f"{index:02d}"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    title = slide.shapes.add_textbox(Inches(0.65), Inches(1.62), Inches(11.2), Inches(1.05))
    tf = title.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = _clean_text(title_text)
    r.font.size = Pt(26 if len(_clean_text(title_text)) > 55 else 30)
    r.font.bold = True; r.font.color.rgb = TEXT

    bullets = bullets or ["Poin utama disesuaikan berdasarkan data calon pelanggan."]
    bullets = [_clean_text(b) for b in bullets if _clean_text(b) and not _is_note(b)][:4]

    positions = [(0.75, 3.00), (6.55, 3.00), (0.75, 4.72), (6.55, 4.72)]
    for i, bullet in enumerate(bullets):
        x, y = positions[i]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.25), Inches(1.28))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = LINE

        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.22), Inches(y+0.30), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background()

        tx = slide.shapes.add_textbox(Inches(x+0.56), Inches(y+0.20), Inches(4.35), Inches(0.88))
        tf = tx.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]; r = p.add_run()
        r.text = bullet
        r.font.size = Pt(14 if len(bullet) > 85 else 15)
        r.font.color.rgb = TEXT

    if notes:
        note_text = _clean_text(notes)
        if note_text:
            note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(6.45), Inches(10.75), Inches(0.40))
            note.fill.solid(); note.fill.fore_color.rgb = RGBColor(236, 253, 245); note.line.color.rgb = RGBColor(187, 247, 208)
            tx = slide.shapes.add_textbox(Inches(0.95), Inches(6.53), Inches(10.3), Inches(0.25))
            p = tx.text_frame.paragraphs[0]; r = p.add_run()
            r.text = f"Narasi: {note_text[:180]}"
            r.font.size = Pt(9); r.font.italic = True; r.font.color.rgb = RGBColor(22, 101, 52)

    add_footer(slide, index=index)

def add_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_top_brand(slide, dark=True)

    a = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.4), Inches(0.6), Inches(3.4), Inches(3.4))
    a.fill.solid(); a.fill.fore_color.rgb = GREEN; a.fill.transparency = 35; a.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.75), Inches(2.15), Inches(9), Inches(0.9))
    p = title.text_frame.paragraphs[0]; r = p.add_run()
    r.text = "Siap lanjut konsultasi?"
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.78), Inches(3.25), Inches(8.3), Inches(1.0))
    tf = sub.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "Zimeira Tech siap menyesuaikan layanan dengan kebutuhan, kondisi, dan prioritas calon pelanggan."
    r.font.size = Pt(18); r.font.color.rgb = RGBColor(203, 213, 225)

    cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(5.05), Inches(4.5), Inches(0.65))
    cta.fill.solid(); cta.fill.fore_color.rgb = GREEN; cta.line.fill.background()
    tf = cta.text_frame; tf.clear()
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "Hubungi Zimeira Tech"
    r.font.bold = True; r.font.size = Pt(18); r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, dark=True)

def create_pptx_from_outline(stamp: str, filename: str, outline_text, project_title="Presentasi Layanan Zimeira Tech"):
    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, project_title)

    for i, s in enumerate(parse_slides(outline_text), start=1):
        add_content_slide(prs, i, s.get("title", f"Slide {i}"), s.get("bullets", []), s.get("notes", ""))

    add_closing_slide(prs)

    path = output_dir / f"{stamp}_{filename}.pptx"
    prs.save(path)
    return str(path)
